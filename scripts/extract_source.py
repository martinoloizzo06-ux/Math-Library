#!/usr/bin/env python3
"""
extract_source.py — Estrae il testo da una fonte (PDF/Word/PPT) in file .md
nella cartella fonti/_estratti/<id-fonte>/, capitolo per capitolo dove possibile.

Uso:
    python3 scripts/extract_source.py <id-fonte>

Esempi:
    python3 scripts/extract_source.py villanacci-successioni
    python3 scripts/extract_source.py axler-ladr
    python3 scripts/extract_source.py hull-ofod

Il catalogo CATALOGO_FONTI.md nella radice del progetto viene usato per
trovare il percorso del file. In alternativa si può passare il percorso
direttamente con --file.

Opzioni:
    --file PATH     Percorso esplicito al file (sovrascrive il catalogo)
    --pages A-B     Estrai solo le pagine A-B (solo PDF; es. --pages 1-30)
    --out DIR       Cartella di output (default: fonti/_estratti/<id>/)
    --split-toc     Tenta di spezzare per capitolo usando il TOC del PDF
"""

import argparse
import re
import sys
from pathlib import Path

# ── Radice del progetto (due livelli sopra scripts/) ─────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
CATALOGO = PROJECT_ROOT / "CATALOGO_FONTI.md"
ESTRATTI = PROJECT_ROOT / "fonti" / "_estratti"


# ── Parsing del catalogo per trovare il percorso ─────────────────────────────

def find_path_in_catalog(source_id: str) -> Path | None:
    if not CATALOGO.exists():
        return None
    text = CATALOGO.read_text(encoding="utf-8")
    # cerca riga "- **Percorso**: `fonti/...`" dopo "### <id>"
    pattern = rf"###\s+{re.escape(source_id)}.*?Percorso.*?`([^`]+)`"
    m = re.search(pattern, text, re.DOTALL)
    if m:
        return PROJECT_ROOT / m.group(1)
    return None


# ── Estrazione PDF ────────────────────────────────────────────────────────────

def extract_pdf(pdf_path: Path, out_dir: Path, page_range: tuple[int, int] | None, split_toc: bool):
    try:
        import pypdf
    except ImportError:
        sys.exit("Installa pypdf: pip install pypdf")

    reader = pypdf.PdfReader(str(pdf_path))
    total = len(reader.pages)
    start, end = (page_range if page_range else (0, total - 1))
    start = max(0, start)
    end = min(total - 1, end)

    print(f"PDF: {pdf_path.name} — {total} pagine totali, estraggo pp. {start+1}–{end+1}")

    # Tenta split per capitolo via outline/TOC
    chapters: list[tuple[str, int, int]] = []  # (titolo, prima_pag, ultima_pag)
    if split_toc:
        try:
            outline = reader.outline
            items = _flatten_outline(reader, outline)
            # filtra solo le pagine nel range
            items = [(t, p) for t, p in items if start <= p <= end]
            for i, (title, pg) in enumerate(items):
                next_pg = items[i + 1][1] - 1 if i + 1 < len(items) else end
                chapters.append((title, pg, next_pg))
        except Exception:
            chapters = []

    if chapters:
        print(f"TOC trovato: {len(chapters)} capitoli/sezioni")
        for idx, (title, pg_start, pg_end) in enumerate(chapters, 1):
            slug = re.sub(r"[^\w\-]", "-", title.lower().strip())[:60].strip("-")
            filename = out_dir / f"{idx:02d}-{slug}.md"
            text = _extract_pages(reader, pg_start, pg_end)
            filename.write_text(f"# {title}\n\n{text}", encoding="utf-8")
            print(f"  [{idx:02d}] pp {pg_start+1}–{pg_end+1} → {filename.name}")
    else:
        # Nessun TOC: split ogni N pagine (default 20) per tenere i file leggibili
        CHUNK = 20
        chunk_idx = 1
        for pg in range(start, end + 1, CHUNK):
            pg_end_chunk = min(pg + CHUNK - 1, end)
            text = _extract_pages(reader, pg, pg_end_chunk)
            filename = out_dir / f"{chunk_idx:02d}-pagine-{pg+1}-{pg_end_chunk+1}.md"
            filename.write_text(text, encoding="utf-8")
            print(f"  [{chunk_idx:02d}] pp {pg+1}–{pg_end_chunk+1} → {filename.name}")
            chunk_idx += 1


def _extract_pages(reader, start: int, end: int) -> str:
    parts = []
    for i in range(start, end + 1):
        txt = reader.pages[i].extract_text() or ""
        # pulizia minima: collassa linee vuote multiple
        txt = re.sub(r"\n{3,}", "\n\n", txt)
        parts.append(f"<!-- pagina {i+1} -->\n{txt}")
    return "\n\n".join(parts)


def _flatten_outline(reader, outline, result=None) -> list[tuple[str, int]]:
    if result is None:
        result = []
    for item in outline:
        if isinstance(item, list):
            _flatten_outline(reader, item, result)
        else:
            try:
                pg = reader.get_destination_page_number(item)
                title = item.title if hasattr(item, "title") else str(item)
                result.append((title, pg))
            except Exception:
                pass
    return result


# ── Estrazione DOCX ───────────────────────────────────────────────────────────

def extract_docx(docx_path: Path, out_dir: Path):
    try:
        from docx import Document
    except ImportError:
        sys.exit("Installa python-docx: pip install python-docx")

    doc = Document(str(docx_path))
    chapters: list[tuple[str, list[str]]] = []
    current_title = "Documento"
    current_paras: list[str] = []

    for para in doc.paragraphs:
        style = para.style.name
        if style.startswith("Heading"):
            if current_paras:
                chapters.append((current_title, current_paras))
            current_title = para.text.strip() or current_title
            current_paras = []
        else:
            if para.text.strip():
                current_paras.append(para.text.strip())

    if current_paras:
        chapters.append((current_title, current_paras))

    print(f"DOCX: {docx_path.name} — {len(chapters)} sezioni")
    for idx, (title, paras) in enumerate(chapters, 1):
        slug = re.sub(r"[^\w\-]", "-", title.lower())[:60].strip("-")
        filename = out_dir / f"{idx:02d}-{slug}.md"
        content = f"# {title}\n\n" + "\n\n".join(paras)
        filename.write_text(content, encoding="utf-8")
        print(f"  [{idx:02d}] {title[:60]} → {filename.name}")


# ── Estrazione PPTX ───────────────────────────────────────────────────────────

def extract_pptx(pptx_path: Path, out_dir: Path):
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("Installa python-pptx: pip install python-pptx")

    prs = Presentation(str(pptx_path))
    print(f"PPTX: {pptx_path.name} — {len(prs.slides)} slide")
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        lines.append(f"## Slide {i}\n\n" + "\n\n".join(slide_texts))

    filename = out_dir / "01-slides.md"
    filename.write_text("\n\n---\n\n".join(lines), encoding="utf-8")
    print(f"  → {filename.name}")


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("source_id", help="ID fonte (es. villanacci-successioni)")
    parser.add_argument("--file", help="Percorso esplicito al file")
    parser.add_argument("--pages", help="Range pagine PDF (es. 1-30)")
    parser.add_argument("--out", help="Cartella di output")
    parser.add_argument("--split-toc", action="store_true", help="Spezza per capitolo via TOC PDF")
    args = parser.parse_args()

    # Risolvi percorso file
    if args.file:
        src = Path(args.file)
    else:
        src = find_path_in_catalog(args.source_id)
        if src is None:
            sys.exit(f"Fonte '{args.source_id}' non trovata in {CATALOGO}. Usa --file per specificare il percorso.")

    if not src.exists():
        sys.exit(f"File non trovato: {src}")

    # Cartella output
    out_dir = Path(args.out) if args.out else ESTRATTI / args.source_id
    out_dir.mkdir(parents=True, exist_ok=True)
    print(f"Output: {out_dir}")

    # Range pagine
    page_range = None
    if args.pages:
        parts = args.pages.split("-")
        page_range = (int(parts[0]) - 1, int(parts[1]) - 1)  # converti a 0-based

    # Dispatch per tipo
    suffix = src.suffix.lower()
    if suffix == ".pdf":
        extract_pdf(src, out_dir, page_range, args.split_toc)
    elif suffix in (".docx", ".doc"):
        extract_docx(src, out_dir)
    elif suffix in (".pptx", ".ppt"):
        extract_pptx(src, out_dir)
    else:
        sys.exit(f"Formato non supportato: {suffix} (supportati: pdf, docx, pptx)")

    print(f"\n✓ Estrazione completata → {out_dir}")


if __name__ == "__main__":
    main()
